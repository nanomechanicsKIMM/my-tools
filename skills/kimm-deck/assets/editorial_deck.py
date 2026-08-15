"""editorial_deck — design.editorial 양식(16:9 설득형 보고 덱) python-pptx 헬퍼.

design.kimm(4:3 보고서 양식)의 자매 시스템. 같은 내용을 '발표해서 설득하는' 자리에 쓴다.
  · 13.333 × 7.5 in · 좌측 라벨 컬럼 + 가로 헤어라인 + 하단 결론 스트립
  · 슬라이드당 메시지 1개 — 헤드라인이 곧 결론
  · 수치는 문장이 아니라 이중 막대 / 대형 스탯 / 실축척 스케일로 보여준다

사용:
    from editorial_deck import *
    prs = new_deck(total=8)
    s = add_slide(prs)
    header(s, 2, "결과", "방법은 재현됐다", "부제 한 줄")
    bar_row(s, 2.2, "지표 A ↓", "0.14 mm", "0.10 mm", 1.0, 0.71, "상회", TEAL)
    closer(s, "한 문장 결론", TEAL)
    foot(s, "덱 이름 · 날짜")
    save(prs, "out.pptx")     # 저장 + 테마 그림자(<p:style>) 제거

전체 규격: design.editorial.md
"""
import os
import re
import zipfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- 팔레트 (design.editorial §2) -----------------------------------------
INK = RGBColor(0x14, 0x16, 0x1A)      # 본문·헤드라인
DEEP = RGBColor(0x10, 0x26, 0x3A)     # 구조색 (표지 룰·기준선)
MUTE = RGBColor(0x8A, 0x90, 0x99)     # 보조 텍스트·축
RULE = RGBColor(0xD8, 0xD4, 0xCC)     # 헤어라인
WASH = RGBColor(0xF3, 0xF1, 0xED)     # 패널 바탕
TEAL = RGBColor(0x0E, 0x7C, 0x86)     # 달성·일치
CORAL = RGBColor(0xD1, 0x49, 0x5B)    # 미달
AMBER = RGBColor(0xE0, 0x9F, 0x3E)    # 진행·주의
GRAYBAR = RGBColor(0xC9, 0xCC, 0xD1)  # 기준(논문·목표) 막대
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TR = "Trebuchet MS"                   # 숫자·라틴 (비번들 폰트 금지)

# ---- 그리드 (design.editorial §1) -----------------------------------------
W, H = 13.333, 7.5
M = 0.75                  # 좌측 여백
R = 12.58                 # 우측 끝
LAB_W = 2.05              # 라벨 컬럼 폭
COL = 3.10                # 콘텐츠 컬럼 시작
COL_W = R - COL           # 9.48
BODY_TOP = 1.78           # 헤어라인 아래 콘텐츠 시작
BODY_BOTTOM = 6.35        # 결론 스트립 위 한계

_TOTAL = [0]              # 페이지 표기 "03 / 08"의 분모


def new_deck(total=0):
    """16:9 빈 덱. total을 주면 헤더에 'nn / total'이 찍힌다."""
    _TOTAL[0] = total
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---- 기본 도형 -------------------------------------------------------------
def box(s, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    return sh


def rule(s, x, y, w, color=RULE, h=0.011):
    """가로 헤어라인 — 이 양식의 유일한 구획 장치."""
    return box(s, x, y, w, h, fill=color)


def text(s, x, y, w, h, runs, size=14, color=INK, bold=False, font=None,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, sp_after=None,
         line_sp=None):
    """runs: str | 줄 리스트 | [(txt, {size,color,bold,font}), ...] 줄 리스트."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [runs]
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if sp_after is not None:
            p.space_after = Pt(sp_after)
        if line_sp is not None:
            p.line_spacing = line_sp
        for txt_, o in (line if isinstance(line, list) else [(line, {})]):
            r = p.add_run(); r.text = txt_
            r.font.size = Pt(o.get("size", size))
            r.font.color.rgb = o.get("color", color)
            r.font.bold = o.get("bold", bold)
            f = o.get("font", font)
            if f:
                r.font.name = f
    return tb


# ---- 페이지 크롬 -----------------------------------------------------------
def header(s, n, kicker, headline, sub=None, size=30):
    """좌측 라벨 컬럼(섹션·페이지) + 헤드라인 + 가로 룰."""
    text(s, M, 0.60, LAB_W, 0.30, kicker, size=13, bold=True, color=TEAL, font=TR)
    label = f"{n:02d} / {_TOTAL[0]:02d}" if _TOTAL[0] else f"{n:02d}"
    text(s, M, 0.90, LAB_W, 0.30, label, size=12, bold=True, color=MUTE, font=TR)
    text(s, COL, 0.48, COL_W, 0.60, headline, size=size, bold=True, color=INK)
    if sub:
        text(s, COL, 1.06, COL_W, 0.34, sub, size=15, color=MUTE)
    rule(s, M, 1.52, R - M)


def closer(s, txt, color=DEEP):
    """하단 결론 스트립 — 굵은 좌측 바 + 한 문장 (design.kimm의 ➡ 배너 대체)."""
    box(s, M, 6.52, 0.09, 0.46, fill=color)
    text(s, M + 0.28, 6.55, R - M - 0.28, 0.42, txt, size=19, bold=True,
         color=color, anchor=MSO_ANCHOR.MIDDLE)


def foot(s, footer_text=""):
    rule(s, M, 7.12, R - M, RULE, 0.008)
    if footer_text:
        text(s, M, 7.19, 8.0, 0.24, footer_text, size=10, color=MUTE, font=TR)


# ---- 콘텐츠 장치 -----------------------------------------------------------
def cover(s, kicker, date, title_lines, subtitle, meta, thesis=None,
          thesis_tail=None):
    """표지 — 좌측 세로 룰 + 대형 제목 + 메타 4열 + 한 줄 주장.

    meta: [(라벨, 값, 부연), ...] 최대 4개.
    thesis / thesis_tail: 마지막 한 문장(뒤쪽 절만 코랄로 강조).
    """
    box(s, M, 1.30, 0.16, 3.20, fill=DEEP)
    text(s, M + 0.45, 1.32, 9.0, 0.34, kicker, size=14, bold=True, color=TEAL, font=TR)
    if date:
        text(s, R - 3.2, 1.32, 3.2, 0.34, date, size=14, bold=True, color=MUTE,
             font=TR, align=PP_ALIGN.RIGHT)
    text(s, M + 0.45, 1.80, 11.0, 1.70, title_lines, size=44, bold=True, color=INK,
         line_sp=1.12)
    text(s, M + 0.45, 3.62, 10.5, 0.50, subtitle, size=17, color=MUTE)
    rule(s, M, 4.30, R - M)
    for i, (k, v1, v2) in enumerate(meta[:4]):
        x = [M, 4.10, 7.45, 10.30][i]
        w = min(3.1, R - x)                      # 마지막 열이 캔버스를 넘지 않게
        text(s, x, 4.52, w, 0.26, k, size=12, bold=True, color=TEAL, font=TR)
        text(s, x, 4.82, w, 0.30, v1, size=16, bold=True, color=INK)
        text(s, x, 5.14, w, 0.28, v2, size=13, color=MUTE)
    rule(s, M, 5.62, R - M)
    if thesis:
        runs = [(thesis, {"color": INK})]
        if thesis_tail:
            runs.append((thesis_tail, {"color": CORAL}))
        text(s, M, 5.86, 11.5, 0.55, [runs], size=24, bold=True)


def stat(s, x, y, w, value, unit, label, note, color=DEEP):
    """대형 숫자 + 단위 + 라벨 + 부연 2줄. 가로 4개까지 배치한다."""
    text(s, x, y, w, 0.80,
         [[(value, {"size": 46, "bold": True, "color": color, "font": TR}),
           (" " + unit, {"size": 17, "bold": True, "color": color, "font": TR})]],
         anchor=MSO_ANCHOR.BOTTOM)
    text(s, x, y + 0.86, w, 0.28, label, size=15, bold=True, color=INK)
    text(s, x, y + 1.16, w, 0.60, note, size=13, color=MUTE, sp_after=2)


def bar_head(s, y=1.78, left="항목", mid="기준 (위) vs 우리 (아래)", right="판정"):
    text(s, M, y, 2.25, 0.26, left, size=12, bold=True, color=MUTE, font=TR)
    text(s, COL, y, 4.0, 0.26, mid, size=12, bold=True, color=MUTE, font=TR)
    if right:
        text(s, R - 1.30, y, 1.30, 0.26, right, size=12, bold=True, color=MUTE,
             font=TR, align=PP_ALIGN.RIGHT)


def bar_row(s, y, name, ref_val, our_val, frac_ref, frac_our, verdict, color,
            x=COL, w=6.05):
    """기준(회색 위) vs 우리(의미색 아래) 이중 막대 한 줄.

    frac_* 는 그 행의 최대값을 1.0으로 둔 정규화 길이 — 행끼리 비교하지 않는다.
    지표 방향이 섞이면 name에 ↓(작을수록 좋음)·↑(클수록 좋음)를 붙일 것.
    """
    text(s, M, y - 0.02, 2.25, 0.30, name, size=14, bold=True, color=INK)
    box(s, x, y + 0.02, w * frac_ref, 0.15, fill=GRAYBAR)
    box(s, x, y + 0.23, w * frac_our, 0.15, fill=color)
    text(s, x + w + 0.14, y - 0.06, 1.35, 0.24, ref_val, size=12, color=MUTE, font=TR)
    text(s, x + w + 0.14, y + 0.19, 1.35, 0.24, our_val, size=12, bold=True,
         color=color, font=TR)
    if verdict:
        text(s, R - 1.30, y + 0.04, 1.30, 0.30, verdict, size=14, bold=True,
             color=color, align=PP_ALIGN.RIGHT)


def scale(s, y, vmin, vmax, ticks, span, band, title=None,
          span_color=CORAL, band_color=DEEP):
    """실축척 증거 그래픽 — 목표 창과 실제 산포를 같은 축에 비율 그대로 그린다.

    span: (lo, hi, 라벨)  — 우리 결과의 범위 (굵은 막대)
    band: (lo, hi, 라벨)  — 목표 창. 폭이 곧 허용오차이므로 임의 두께를 쓰지 말 것.
    '허용오차의 N배'라는 문장을 눈으로 확인시키는 것이 이 장치의 목적이다.
    """
    def px(v):
        return M + (R - M) * (v - vmin) / float(vmax - vmin)

    if title:
        text(s, M, y, 9.0, 0.30, title, size=15, bold=True, color=INK)
    axis_y = y + 1.36
    box(s, px(band[0]), y + 0.65, max(px(band[1]) - px(band[0]), 0.012), 0.71,
        fill=band_color)
    box(s, px(span[0]), y + 0.82, px(span[1]) - px(span[0]), 0.34, fill=span_color)
    text(s, px(span[0]) + 0.55, y + 0.46, 5.0, 0.26, span[2], size=12, bold=True,
         color=span_color, font=TR)
    rule(s, M, axis_y, R - M, MUTE, 0.014)
    for v in ticks:
        box(s, px(v), axis_y, 0.011, 0.12, fill=MUTE)
        text(s, px(v) - 0.35, axis_y + 0.16, 0.70, 0.24, str(v), size=11, color=MUTE,
             font=TR, align=PP_ALIGN.CENTER)
    mid = (band[0] + band[1]) / 2.0
    text(s, px(mid) - 1.90, axis_y + 0.50, 3.8, 0.26, "▲ " + band[2], size=12,
         bold=True, color=band_color, font=TR, align=PP_ALIGN.CENTER)


def card(s, x, y, w, h, color, num, title_, verdict, lines, next_line=None):
    """진단 카드 — 번호 · 제목 · 판정 · 근거 · (선택) 다음 수."""
    box(s, x, y, w, h, fill=WHITE, line=RULE)
    box(s, x, y, w, 0.09, fill=color)
    text(s, x + 0.25, y + 0.22, 0.6, 0.40, num, size=24, bold=True, color=color, font=TR)
    text(s, x + 0.25, y + 0.74, w - 0.45, 0.32, title_, size=18, bold=True, color=INK)
    text(s, x + 0.25, y + 1.12, w - 0.45, 0.28, verdict, size=13, bold=True, color=color)
    rule(s, x + 0.25, y + 1.52, w - 0.50)
    # 본문 높이는 카드 높이에서 역산한다 — '다음 수' 박스와 겹치지 않게
    body_top = y + 1.68
    next_top = y + h - 0.20 - 0.62
    body_h = max((next_top if next_line else y + h - 0.20) - body_top - 0.10, 0.30)
    text(s, x + 0.25, body_top, w - 0.45, body_h, lines, size=13, color=INK, sp_after=3)
    if next_line:
        box(s, x + 0.25, next_top, w - 0.50, 0.62, fill=WASH)
        text(s, x + 0.40, next_top + 0.10, w - 0.80, 0.42, next_line, size=12,
             bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)


def timeline(s, items, x=M, y=2.00, step=0.90, when_w=1.55):
    """수직 타임라인 — items: [(시점, 색, 제목, 설명), ...]"""
    tx = x + 0.70 + when_w + 0.25
    avail = max(R - tx, 1.0)                     # 우측 마진을 넘지 않게
    box(s, x + 0.30, y - 0.05, 0.03, step * len(items) - 0.05, fill=RULE)
    for i, (when, c, ttl, desc) in enumerate(items):
        yy = y + step * i
        box(s, x + 0.20, yy + 0.06, 0.24, 0.24, fill=c, shape=MSO_SHAPE.OVAL)
        text(s, x + 0.70, yy, when_w, 0.30, when, size=15, bold=True, color=c, font=TR)
        text(s, tx, yy - 0.02, min(3.0, avail), 0.30, ttl, size=16, bold=True,
             color=INK)
        text(s, tx, yy + 0.32, avail, 0.34, desc, size=13, color=MUTE)


def panel(s, y, title_, lines, x=M, w=None, h=1.62):
    """WASH 바탕 강조 패널 (규율·요약 등)."""
    w = w if w is not None else R - M
    box(s, x, y, w, h, fill=WASH)
    text(s, x + 0.30, y + 0.23, w - 0.60, 0.30, title_, size=16, bold=True, color=DEEP)
    text(s, x + 0.30, y + 0.63, w - 0.60, h - 0.80, lines, size=14, color=INK,
         sp_after=5)


def save(prs, out):
    """저장 후 <p:style> 제거 — python-pptx 도형의 테마 그림자를 없앤다."""
    prs.save(out)
    tmp = out + ".tmp"
    with zipfile.ZipFile(out) as zin, \
            zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for it in zin.infolist():
            data = zin.read(it.filename)
            if it.filename.startswith("ppt/slides/") and it.filename.endswith(".xml"):
                data = re.sub(r"<p:style>.*?</p:style>", "",
                              data.decode("utf-8")).encode("utf-8")
            zout.writestr(it, data)
    os.replace(tmp, out)
    return out


def audit(path):
    """구조 감사 — 캔버스 이탈 · 테마 그림자 · 비번들 폰트. QA 필수 단계."""
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    out = []
    for i, s in enumerate(prs.slides, 1):
        for shp in s.shapes:
            if shp.left is None:
                continue
            r, b = shp.left + (shp.width or 0), shp.top + (shp.height or 0)
            if shp.left < 0 or shp.top < 0 or r > sw or b > sh:
                out.append(f"slide {i}: {shp.name} 캔버스 이탈")
    with zipfile.ZipFile(path) as z:
        names = [n for n in z.namelist()
                 if n.startswith("ppt/slides/") and n.endswith(".xml")]
        for n in names:
            raw = z.read(n)
            if b"<p:style>" in raw:
                out.append(f"{n}: 테마 그림자(<p:style>) 잔여")
            for f in set(re.findall(rb'typeface="([^"]+)"', raw)):
                if f.decode() not in (TR, "Arial", "Calibri", "Georgia"):
                    out.append(f"{n}: 비번들 폰트 {f.decode()}")
    return out
