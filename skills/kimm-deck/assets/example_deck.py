"""example_deck.py — kimm-deck 스킬 검증 예제 (patent-incubation-auto 사용 안내 7슬라이드).

design.kimm 아키타입 A·B·C·D·E·G·H를 kimm_deck 헬퍼로 구현한 레퍼런스.
실행:  python3 example_deck.py [출력경로]
"""
import os
import sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from kimm_deck import *

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(os.getcwd(), "example-kimm-deck.pptx")
FOOT = "patent-incubation-auto quick guide"

prs = new_deck()

def slide():
    return add_slide(prs)

# ===========================================================================
# S1 — 아키타입 A: 표지
s = slide()
s.shapes.add_picture(os.path.join(HERE, "bg_title.png"), 0, 0, Inches(10), Inches(7.5))
text(s, 0.03, 0.02, 9.6, 0.44, "KIMM Patent Incubation Toolkit", size=20, bold=True,
     color=WHITE, font=TREB)
text(s, 0.2, 1.2, 9.6, 1.7,
     ["patent-incubation-auto — 사용 방법 안내",
      [("TRIZ 분석 · 선행특허 조사 · 발명내용설명서 HWPX 자동 생성", {"size": 20})]],
     size=34, bold=True, color=WHITE, font=TREB, align=PP_ALIGN.CENTER)
d = box(s, 1.9, 3.37, 6.3, 0.55, fill=DATEBAR)
text(s, 1.9, 3.40, 6.3, 0.50, "August 8th, 2026", size=20, bold=True, color=WHITE,
     font=TREB, align=PP_ALIGN.CENTER)
p = box(s, 1.2, 4.2, 7.6, 2.7, fill=PANEL)
rows = [("입력", "기술분야 · 해결 과제 · 핵심 아이디어 (3요소)"),
        ("분석", "TRIZ 모순/IFR + KIPRIS · 국제 선행조사"),
        ("검증", "청구항 하드닝 + Critic 이중 게이트"),
        ("출력", "KIMM 직무발명내용설명서 HWPX + 도면"),
        ("실행", "\"/patent-incubation-auto\" 또는 \"발명내용설명서 써줘\"")]
text(s, 1.5, 4.35, 7.0, 2.45,
     [[(k + "        ", {"bold": True}), (v, {})] for k, v in rows],
     size=17, color=WHITE, font=TREB, sp_after=8)

# ===========================================================================
# S2 — 아키타입 E: 파이프라인(타임라인) + 준비물 카드
s = slide()
chrome(s, "Workflow", "Pipeline — Step 0 to 8", 2, footer_text=FOOT)
text(s, 0.35, 1.18, 3.0, 0.4, "Steps", size=20, bold=True, color=NAVY, font=TREB)
box(s, 0.92, 1.66, 0.03, 4.15, fill=TLGRAY)
steps = [("Step 0", NAVY, "입력 수집 — 아이디어 3요소 + 발명자/기관"),
         ("Step 1–2", NAVY, "TRIZ 시스템 분석 · 모순/IFR 도출"),
         ("Step 4–5.5", AMBER, "정량 평가 · 선행조사 · 특허성 재채점"),
         ("Step 6–6.5", BLUE, "발명내용설명서 작성 + 청구항 하드닝"),
         ("Step 6b–6e", BLUE, "도면 생성 · 인용 게이트 · Critic ×2"),
         ("Step 7–8", GREEN, "HWPX 변환 · 최종 안내")]
for i, (d_, c, desc) in enumerate(steps):
    y = 1.71 + 0.72 * i
    dot(s, 0.85, y, c)
    text(s, 1.15, y - 0.12, 3.9, 0.35, d_, size=16, bold=True, color=c, font=TREB)
    text(s, 1.15, y + 0.16, 3.9, 0.40, desc, size=14, color=INK)
text(s, 5.17, 1.18, 4.4, 0.4, "준비물", size=16, bold=True, color=NAVY, font=TREB)
cards = [(AMBER, "1", "아이디어 3요소", "기술분야 · 해결 과제 · 핵심 아이디어. md 문서가 있으면 자동 감지."),
         (BLUE, "2", "KIPRIS API 키", "~/Claude_Work/.env. 없으면 선행조사만 degraded 모드로 진행."),
         (GREEN, "3", "hwpx 스킬", "HWPX 변환 엔진. my-tools setup 시 자동 설치.")]
for k, (c, num, t, b) in enumerate(cards):
    y = 1.56 + 1.44 * k
    box(s, 5.17, y, 4.48, 1.24, fill=WHITE, line=BORDER)
    box(s, 5.17, y, 0.52, 1.24, fill=c)
    text(s, 5.17, y, 0.52, 1.24, num, size=22, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.82, y + 0.10, 3.7, 0.35, t, size=16, bold=True, color=INK)
    text(s, 5.82, y + 0.48, 3.7, 0.70, b, size=13, color=INK)
banner(s, "아이디어 세 줄이면 파이프라인이 신고서까지 간다")

# ===========================================================================
# S3 — 아키타입 B: 산출물 표 + 우측 담당 Phase 라벨
s = slide()
chrome(s, "Outputs", "Main outputs (output/ 디렉터리)", 3, footer_text=FOOT)
data = [("발명내용설명서.md", "MD", "§1–§9 + 부록 A(TRIZ) · B(회피설계)", "Obsidian"),
        ("prior_art.json", "JSON", "선행특허 · 자기공지(grace 기한 포함)", "KIPRIS"),
        ("diagrams/*.png", "PNG", "600 dpi 도면 (부호 없음 — 이름 라벨만)", "≥ 5매"),
        ("figures_deck.pptx", "PPTX", "편집 가능 도면 덱 (SVG 벡터 삽입)", "편집용"),
        ("발명내용설명서.hwpx", "HWPX", "KIMM 직무발명신고서 최종본", "제출용"),
        ("manifest.json", "JSON", "Phase별 진행 상태 기록", "재개용")]
tw = 8.78
tbl_sh = s.shapes.add_table(len(data) + 1, 4, Inches(0.09), Inches(1.10),
                            Inches(tw), Inches(0.55 + 0.36 * len(data)))
t = tbl_sh.table
t.columns[0].width = Inches(2.3); t.columns[1].width = Inches(0.9)
t.columns[2].width = Inches(4.2); t.columns[3].width = Inches(1.38)
hdr = ["산출물", "형식", "내용", "비고"]
t.rows[0].height = Inches(0.55)
for j, htxt in enumerate(hdr):
    c = t.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = BLUE
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    pgh = c.text_frame.paragraphs[0]; pgh.alignment = PP_ALIGN.CENTER
    r = pgh.add_run(); r.text = htxt
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
for i, row in enumerate(data, 1):
    t.rows[i].height = Inches(0.36)
    for j, val in enumerate(row):
        c = t.cell(i, j); c.fill.solid()
        c.fill.fore_color.rgb = BAND if i % 2 == 0 else WHITE
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        pgh = c.text_frame.paragraphs[0]
        pgh.alignment = PP_ALIGN.LEFT if j in (0, 2) else PP_ALIGN.CENTER
        r = pgh.add_run(); r.text = val
        r.font.size = Pt(14); r.font.color.rgb = INK
        if j == 0: r.font.bold = True
owners = [("Phase 6", 1), ("Phase 5", 2), ("Phase 6b", 3), ("Phase 6b", 4),
          ("Phase 7", 5), ("orchestr.", 6)]
for lab, ri in owners:
    y = 1.10 + 0.55 + 0.36 * (ri - 1) + 0.03
    text(s, 8.92, y, 1.05, 0.3, lab, size=13, bold=True, color=NAVY, wrap=False)
banner(s, "HWPX 최종본과 함께 모든 중간 산출물이 남는다")

# ===========================================================================
# S4 — 아키타입 D: 스탯 카드 2×2 (품질 게이트)
s = slide()
chrome(s, "Features", "Quality gates & guarantees", 4, footer_text=FOOT)
stats = [(0.35, 1.22, GREEN, "15", "청구항 하드닝 점검 항목",
          ["112(b) · 트리 정합 · 회피설계 매트릭스", "발견 즉시 자동 수정 후 보고"]),
         (5.17, 1.22, BLUE, "×2", "Critic 이중 게이트 (opus)",
          ["등록가능성: PASS / FIX / BLOCK", "사업화: 회피설계 · 침해 입증 관점"]),
         (0.35, 3.70, BLUE, "≥5", "자동 생성 컬러 도면",
          ["SVG → 편집 가능 PPTX → 600 dpi PNG", "도면 부호 없음 — 변리사 작업과 무충돌"]),
         (5.17, 3.70, ORANGE, "12개월", "자기공지 유예 기한 자동 경고",
          ["논문 선공개 발견 시 grace_deadline 기록", "KIMM 최다 무효 사유를 사전 차단"])]
for x, y, c, num, lab, subs in stats:
    box(s, x, y, 4.48, 2.28, fill=WHITE, line=BORDER)
    box(s, x, y, 4.48, 0.07, fill=c)
    text(s, x + 0.22, y + 0.22, 4.0, 0.75, num, size=40, bold=True, color=c)
    text(s, x + 0.22, y + 1.02, 4.1, 0.35, lab, size=16, bold=True, color=INK)
    text(s, x + 0.22, y + 1.40, 4.1, 0.80, subs, size=13, color=INK)
banner(s, "점검은 자동, 판단 근거는 전부 문서에 남긴다")

# ===========================================================================
# S5 — 아키타입 C: 요약 2박스 (입력 방법 A/B)
s = slide()
chrome(s, "Usage", "How to run — 두 가지 입력 방법", 5, footer_text=FOOT)
boxes = [(0.35, "입력 방법 A   (대화형, 기본)",
          [[("실행 :", {"bold": True}), ("  \"/patent-incubation-auto\"", {})],
           "",
           [("3요소 입력", {"bold": True})],
           "   기술분야 — 예: 마이크로LED 제조",
           "   해결 과제 — 예: 인터포저 비용·시간 절감",
           "   핵심 아이디어 — 예: 레이저 직접 전사",
           "",
           [("옵션", {"bold": True})],
           "   발명자명(공동발명자 쉼표 구분) · 소속기관",
           "   출력 디렉토리 (기본: ./output/)"]),
         (5.17, "입력 방법 B   (문서 기반, 자동 감지)",
          [[("실행 :", {"bold": True}), ("  \"현재 폴더의 md 파일을 토대로", {})],
           "           발명내용설명서 써줘\"",
           "",
           [("동작", {"bold": True})],
           "   현재 디렉토리 .md 자동 탐색·분석",
           "   3요소를 문서에서 추출 후 확인 질문",
           "",
           [("적합", {"bold": True})],
           "   아이디어 메모·실험 노트가 이미 있을 때",
           "   LLM_wiki 노트 기반 발명 신고"])]
for x, hdr_txt, body in boxes:
    hb = box(s, x, 1.22, 4.48, 0.42, fill=BLACK)
    text(s, x + 0.12, 1.26, 4.3, 0.35, hdr_txt, size=16, bold=True, color=WHITE)
    box(s, x, 1.64, 4.48, 4.20, fill=WHITE, line=BORDER)
    text(s, x + 0.15, 1.80, 4.2, 3.9, body, size=15, color=INK, sp_after=2)
banner(s, "아이디어 3요소만 넣으면 HWPX 신고서까지 자동")

# ===========================================================================
# S6 — 아키타입 G: 3열 상태 보드 + 주의 박스
s = slide()
chrome(s, "Guide", "적합 · 주의 · 대안", 6, footer_text=FOOT)
cols = [(0.35, GREEN, "이럴 때 적합", ["아이디어가 이미 구체적일 때", "무인 완전 자동 생성을 원할 때",
         "직무발명 신고 마감이 급할 때", "md 메모에서 바로 신고서로"]),
        (3.47, AMBER, "주의  (degraded)", ["KIPRIS 키 없음 → 선행조사 축소,", "  §3·§4·§8 수동 보완 필요",
         "논문 선공개 존재 → 12개월", "  공지예외 기한 경고 확인"]),
        (6.59, GRAY, "다른 스킬이 맞는 경우", ["아이디어가 유동적 →", "  patent-incubation-interactive",
         "출원 초안 검토 → draft-review", "거절 대응 → patent-defence"])]
for x, c, hdr_txt, items in cols:
    box(s, x, 1.22, 2.98, 0.42, fill=c)
    text(s, x + 0.10, 1.26, 2.8, 0.35, hdr_txt, size=15, bold=True, color=WHITE)
    box(s, x, 1.64, 2.98, 3.13, fill=WHITE, line=BORDER)
    text(s, x + 0.12, 1.78, 2.76, 2.9, items, size=14, color=INK, sp_after=3)
box(s, 0.35, 5.02, 9.30, 1.18, fill=CREAM, line=AMBER)
text(s, 0.55, 5.12, 8.9, 0.32, "TRIZ 흔적은 최종 문서에 남지 않는다", size=16, bold=True, color=AMBER)
text(s, 0.55, 5.46, 8.9, 0.66,
     "§1–§9와 HWPX에는 TRIZ·IFR·원리 번호가 일반 기술 용어로 변환되어 서술되고, "
     "분석 과정 원본은 부록 A에만 기록된다. 도면 부호·[도 N] 번호도 사용하지 않는다 — 출원용은 변리사가 별도 작성.",
     size=13, color=INK)
banner(s, "막히면 interactive 모드로 전환해 단계별로 진행")

# ===========================================================================
# S7 — 아키타입 H: 클로징
s = slide()
s.shapes.add_picture(os.path.join(HERE, "bg_close.png"), 0, Inches(1.0), Inches(10), Inches(4.5))
pn = box(s, 0.5, 2.55, 9.0, 1.55, fill=PANEL)
text(s, 0.5, 2.75, 9.0, 0.7, "아이디어 세 줄로 시작하세요.", size=30, bold=True,
     color=WHITE, font=TREB, align=PP_ALIGN.CENTER)
text(s, 0.5, 3.50, 9.0, 0.4, "patent-incubation-auto  ·  design.kimm sample  ·  2026-08-08",
     size=16, color=SUBGRAY, font=TREB, align=PP_ALIGN.CENTER)
footer(s, 7, FOOT)

print("saved:", save(prs, OUT))
