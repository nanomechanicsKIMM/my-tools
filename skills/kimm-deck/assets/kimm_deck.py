"""kimm_deck — design.kimm 양식(4:3 KIMM 보고 덱) python-pptx 헬퍼.

사용:
    from kimm_deck import *
    prs = new_deck()
    s = add_slide(prs)
    chrome(s, "Outputs", "슬라이드 제목", n=2, footer_text="deck footer")
    banner(s, "한 문장 결론")
    save(prs, "out.pptx")          # 저장 + 테마 그림자(<p:style>) 제거

전체 규격: design.kimm.md (아키타입 A~H 좌표 포함).
"""
import re
import zipfile
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- 색상 토큰 (design.kimm §2) -------------------------------------------
NAVY = RGBColor(0x1F, 0x38, 0x64)      # 구조 프라이머리
PANEL = RGBColor(0x0B, 0x10, 0x30)     # 다크 패널
DATEBAR = RGBColor(0x03, 0x03, 0x1B)   # 표지 날짜바
INK = RGBColor(0x26, 0x26, 0x26)       # 본문
GREEN = RGBColor(0x2E, 0x7D, 0x32)     # 달성·완료
AMBER = RGBColor(0xB5, 0x6E, 0x00)     # 격차·준비 중·주의
BLUE = RGBColor(0x5B, 0x9B, 0xD5)      # 중간·예정, 표 헤더
ORANGE = RGBColor(0xED, 0x7D, 0x31)    # 미달
RED = RGBColor(0xC0, 0x00, 0x00)       # NA·경고
GRAY = RGBColor(0x80, 0x80, 0x80)      # 중립
BORDER = RGBColor(0x40, 0x40, 0x40)    # 카드 테두리
BAND = RGBColor(0xDE, 0xEA, 0xF6)      # 표 밴드
CREAM = RGBColor(0xFD, 0xF2, 0xDF)     # 주의 박스 배경
TLGRAY = RGBColor(0xC0, 0xC0, 0xC0)    # 타임라인 룰
SUBGRAY = RGBColor(0x9F, 0xB0, 0xC0)   # 클로징 부제
RULE = RGBColor(0x7F, 0x7F, 0x7F)      # 제목 헤어라인
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

TREB = "Trebuchet MS"                  # 제목·칩·배너 (HY·나눔 폰트 금지)


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(10)       # 4:3 (design.kimm §1)
    prs.slide_height = Inches(7.5)
    return prs


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    return sh


def dot(s, x, y, color, d=0.16):
    """타임라인 원형 도트 (design.kimm §5-E)."""
    return box(s, x, y, d, d, fill=color, shape=MSO_SHAPE.OVAL)


def text(s, x, y, w, h, runs, size=14, color=INK, bold=False, font=None,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, sp_after=None):
    """runs: str | 줄 리스트 | [(txt, {size,color,bold,font}), ...] 줄 리스트."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    if isinstance(runs, str):
        runs = [runs]
    first = True
    for line in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if sp_after is not None:
            p.space_after = Pt(sp_after)
        segs = line if isinstance(line, list) else [(line, {})]
        for txt_, o in segs:
            r = p.add_run(); r.text = txt_
            r.font.size = Pt(o.get("size", size))
            r.font.color.rgb = o.get("color", color)
            r.font.bold = o.get("bold", bold)
            f = o.get("font", font)
            if f:
                r.font.name = f
    return tb


def chip(s, label):
    """카테고리 칩 (design.kimm §4-1)."""
    c = box(s, 0, 0, 1.98, 0.40, fill=NAVY)
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = TREB
    return c


def title(s, txt, size=31):
    """제목 + 헤어라인 룰 (design.kimm §4-2)."""
    text(s, 0.25, 0.35, 9.4, 0.62, txt, size=size, bold=True, color=BLACK, font=TREB)
    box(s, 0.25, 1.00, 9.50, 0.014, fill=RULE)


def banner(s, txt):
    """➡ 결론 배너 (design.kimm §4-3)."""
    box(s, 0.62, 6.47, 8.26, 0.02, fill=BLACK)
    box(s, 0.62, 7.05, 8.26, 0.02, fill=BLACK)
    text(s, 0.62, 6.55, 8.26, 0.44, "➡ " + txt, size=20, bold=True,
         color=NAVY, font=TREB, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(s, n, footer_text=""):
    """푸터 (design.kimm §4-4 근사 사양)."""
    box(s, 0, 7.28, 10, 0.22, fill=RGBColor(0xA6, 0xA6, 0xA6))
    if footer_text:
        text(s, 0.25, 7.25, 7.5, 0.28, footer_text, size=11, bold=True,
             color=RGBColor(0x59, 0x59, 0x59))
    box(s, 8.85, 7.25, 1.15, 0.28, fill=PANEL)
    text(s, 8.85, 7.25, 1.15, 0.28, str(n), size=11, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)


def chrome(s, label, ttl, n, size=31, footer_text=""):
    chip(s, label); title(s, ttl, size); footer(s, n, footer_text)


def save(prs, out):
    """저장 후 <p:style> 제거 — 테마 그림자 방지 (design.kimm §6)."""
    prs.save(out)
    tmp = out + ".tmp"
    with zipfile.ZipFile(out) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for it in zin.infolist():
            data = zin.read(it.filename)
            if it.filename.startswith("ppt/slides/") and it.filename.endswith(".xml"):
                data = re.sub(r"<p:style>.*?</p:style>", "",
                              data.decode("utf-8")).encode("utf-8")
            zout.writestr(it, data)
    os.replace(tmp, out)
    return out
