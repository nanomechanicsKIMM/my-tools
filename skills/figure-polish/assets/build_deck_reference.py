# -*- coding: utf-8 -*-
"""
Assemble the full 12-figure deck (figures_v2/display_figures.pptx).

Slide plan:
  1  fig1  native schematic            7  fig7  native schematic (curve a + b/c/d)
  2  fig2  image (complex data)        8  fig8  image (complex data)
  3  fig3  native table                9  fig9  image (complex data)
  4  fig4  native schematic           10  fig10 native schematic + shape bar chart
  5  fig5  native schematic + table   11  fig11 image (complex data)
  6  fig6  native schematic           12  fig12 native schematic
"""
import os
from pptx.util import Mm
import pptx_figkit as K
import fig1_slide, fig4_slide, fig5_slide, fig6_slide, fig10_slide, fig12_slide

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "figures_v2")
GRAY = "7F7F7F"


def image_slide(prs, name):
    path = os.path.join(IMG, name + "_600dpi.png")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    margin = 8.0
    maxw = K.SLIDE_W - 2 * margin
    maxh = K.SLIDE_H - 2 * margin
    pic = slide.shapes.add_picture(path, Mm(margin), Mm(margin), width=Mm(maxw))
    if pic.height > Mm(maxh):
        pic.width = int(pic.width * (Mm(maxh) / pic.height))
        pic.height = Mm(maxh)
    pic.left = int((prs.slide_width - pic.width) / 2)
    pic.top = int((prs.slide_height - pic.height) / 2)
    return slide


def fig3_table(prs):
    s = K.Slide(prs, 120, 60)
    data = [
        ["항목", "LCD", "OLEDoS", "Micro-LED", "QD-EL"],
        ["명암비", "5,000:1", "∞", "∞", "∞"],
        ["발광효율 (lm/W)", "~65 (중)", "~20 (하)", "~100 (상)", "하"],
        ["휘도 (cd/m^{2})", "3,000", "1,500", "100,000", "~100"],
        ["수명", "5~8년", "~4년", ">10년", "<10^{5} 시간"],
        ["응답시간", "5 ms", "110 µs", "~0.2 ns", "LED와 유사"],
    ]
    tbl = s.table(0.0, 0.12, 1.0, 0.97, data,
                  col_w=[0.30, 0.175, 0.175, 0.175, 0.175], fs=8.5)
    # Micro-LED column highlight (col index 3)
    for i in range(len(data)):
        cell = tbl.cell(i, 3)
        cell.fill.solid()
        cell.fill.fore_color.rgb = K.C("21708A" if i == 0 else "D7EBEF")
    # bold vermilion: Micro-LED 휘도 & 응답시간
    for (i, j) in ((3, 3), (5, 3)):
        for r in tbl.cell(i, j).text_frame.paragraphs[0].runs:
            r.font.bold = True
            r.font.color.rgb = K.C("D55E00")
    s.textbox_frac(0.0, 0.0, 1.0, 0.10,
                   "데이터: 저자 벤치마크 자료 (Adv. Funct. Mater. 2019, adfm.201808075 기반)",
                   7.5, GRAY, align="left")
    return s


def main():
    prs = K.new_deck()
    fig1_slide.add_slide(prs, IMG)          # 1
    image_slide(prs, "figure2_requirements")  # 2
    fig3_table(prs)                          # 3
    fig4_slide.add_slide(prs, IMG)           # 4
    fig5_slide.add_slide(prs, IMG)           # 5
    fig6_slide.add_slide(prs, IMG)           # 6
    image_slide(prs, "figure7_transfer_mechanics")  # 7 (curve panel = complex -> image)
    image_slide(prs, "figure8_sizeeffect")   # 8
    image_slide(prs, "figure9_red_comparison")  # 9
    fig10_slide.add_slide(prs, IMG)          # 10
    image_slide(prs, "figure11_landscape")   # 11
    fig12_slide.add_slide(prs, IMG)          # 12
    out = os.path.join(IMG, "display_figures.pptx")
    prs.save(out)
    print("saved:", out)
    for i, sl in enumerate(prs.slides, 1):
        print(f"  slide {i:2d}: {len(sl.shapes)} shapes")


if __name__ == "__main__":
    main()
