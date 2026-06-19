# -*- coding: utf-8 -*-
"""
pptx_figkit — shared toolkit to rebuild matplotlib figures as native, editable
PowerPoint slides (shapes + text boxes), one slide per figure in a 16:9 deck.

Design: each figure is mapped from its matplotlib coordinate space onto the slide
with a single uniform "contain" scale, so proportions and relative font sizes
match the verified PNG. Fonts are scaled by the same factor (pt * SP).

Usage (one module per figure):
    import pptx_figkit as K
    def add_slide(prs, img_dir):
        s = K.Slide(prs, fig_w_mm=183, fig_h_mm=72)        # native figure size
        ax = s.axes(0.005, 0.995, 0.02, 0.98, xlim=(0,18.3), ylim=(0,7.2))
        ax.rect(1.6, 4.6, 6.0, 0.6, "D7E3F0", ec="1A1A1A", label="...", label_fs=7.5)
        ax.text(4.6, 6.95, "제목", 9.3, "1F3A5F", ha="center", bold=True)
        ...
Coordinates mirror matplotlib: rect(x, y, w, h) uses bottom-left; y grows up.
Superscripts: write ^{...}  ->  real PowerPoint superscript run.
"""
import re
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

SLIDE_W = 338.667          # mm (16:9)
SLIDE_H = 190.5
MARGIN = 8.0
FONT = "Pretendard"
INK = "1A1A1A"


def C(h):
    return RGBColor.from_string(h.lstrip("#").upper())


def new_deck():
    prs = Presentation()
    prs.slide_width = Mm(SLIDE_W)
    prs.slide_height = Mm(SLIDE_H)
    return prs


# --------------------------------------------------------------- text runs
def _font(run, fs_pt, color, bold=False, italic=False):
    f = run.font
    f.size = Pt(fs_pt); f.bold = bold; f.italic = italic; f.name = FONT
    f.color.rgb = C(color)
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", FONT)
    return rPr


def _add_runs(p, text, fs_pt, color, bold=False, italic=False):
    for part in re.split(r"(\^\{[^}]*\})", str(text)):
        if not part:
            continue
        if part.startswith("^{") and part.endswith("}"):
            r = p.add_run(); r.text = part[2:-1]
            rPr = _font(r, fs_pt, color, bold, italic); rPr.set("baseline", "30000")
        else:
            r = p.add_run(); r.text = part
            _font(r, fs_pt, color, bold, italic)


def _fill_tf(tf, lines, fs_pt, color, bold, italic, align, anchor, line_spacing):
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    if isinstance(lines, str):
        lines = lines.split("\n")
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        _add_runs(p, ln, fs_pt, color, bold, italic)


def _no_shadow(shp):
    try:
        shp.shadow.inherit = False
    except Exception:
        pass


_ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


class Slide:
    def __init__(self, prs, fig_w_mm, fig_h_mm, fit="contain"):
        self.prs = prs
        self.slide = prs.slides.add_slide(prs.slide_layouts[6])
        avail_w, avail_h = SLIDE_W - 2 * MARGIN, SLIDE_H - 2 * MARGIN
        if fit == "width":
            self.SP = avail_w / fig_w_mm
        else:
            self.SP = min(avail_w / fig_w_mm, avail_h / fig_h_mm)
        self.FW = fig_w_mm * self.SP
        self.FH = fig_h_mm * self.SP
        self.X0 = (SLIDE_W - self.FW) / 2.0
        self.Y0 = (SLIDE_H - self.FH) / 2.0

    # figure-fraction (fy from bottom) -> mm top-origin
    def fig_xy(self, fx, fy):
        return self.X0 + fx * self.FW, self.Y0 + (1.0 - fy) * self.FH

    def pt(self, x):
        return x * self.SP            # scaled font/line points

    def axes(self, left, right, bottom, top, xlim=(0, 1), ylim=(0, 1),
             xlog=False, ylog=False):
        return Ax(self, left, right, bottom, top, xlim, ylim, xlog, ylog)

    # ---- figure-level objects -------------------------------------------
    def image(self, path, fx0=0.0, fy0=0.0, fx1=1.0, fy1=1.0):
        l, t = self.fig_xy(fx0, fy1)
        r, b = self.fig_xy(fx1, fy0)
        self.slide.shapes.add_picture(path, Mm(l), Mm(t), Mm(r - l), Mm(b - t))

    def textbox_frac(self, fx, fy, fw, fh, lines, fs, color, bold=False,
                     italic=False, align="center", anchor=MSO_ANCHOR.MIDDLE,
                     line_spacing=None):
        l, t = self.fig_xy(fx, fy + fh)
        tb = self.slide.shapes.add_textbox(Mm(l), Mm(t), Mm(fw * self.FW), Mm(fh * self.FH))
        _fill_tf(tb.text_frame, lines, self.pt(fs), color, bold, italic,
                 _ALIGN[align], anchor, line_spacing)
        return tb

    def banner(self, fx0, fx1, fy_lo, fy_hi, lines, fs, fc, color="FFFFFF",
               round_adj=0.5):
        l, t = self.fig_xy(fx0, fy_hi)
        r, b = self.fig_xy(fx1, fy_lo)
        shp = self.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Mm(l), Mm(t), Mm(r - l), Mm(b - t))
        _no_shadow(shp); shp.fill.solid(); shp.fill.fore_color.rgb = C(fc)
        shp.line.fill.background()
        try:
            shp.adjustments[0] = round_adj
        except Exception:
            pass
        tf = shp.text_frame
        _fill_tf(tf, lines, self.pt(fs), color, False, False, PP_ALIGN.CENTER,
                 MSO_ANCHOR.MIDDLE, None)
        tf.margin_left = Mm(2); tf.margin_right = Mm(2)
        return shp

    def table(self, fx0, fy0, fx1, fy1, data, col_w=None, header_fc="1F3A5F",
              header_color="FFFFFF", body_color="1A1A1A", fs=8.0,
              row_band="F2F4F7", align="center"):
        """data: list of rows; row0 = header."""
        l, t = self.fig_xy(fx0, fy1)
        r, b = self.fig_xy(fx1, fy0)
        nr, nc = len(data), len(data[0])
        gframe = self.slide.shapes.add_table(nr, nc, Mm(l), Mm(t), Mm(r - l), Mm(b - t))
        tbl = gframe.table
        tbl.first_row = False; tbl.horz_banding = False
        if col_w:
            total = sum(col_w)
            for j, w in enumerate(col_w):
                tbl.columns[j].width = Mm((r - l) * w / total)
        for i, row in enumerate(data):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.margin_left = Mm(1.5); cell.margin_right = Mm(1.5)
                cell.margin_top = Mm(0.5); cell.margin_bottom = Mm(0.5)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.fill.solid()
                if i == 0:
                    cell.fill.fore_color.rgb = C(header_fc); col = header_color; bold = True
                else:
                    cell.fill.fore_color.rgb = C("FFFFFF" if i % 2 else row_band)
                    col = body_color; bold = (j == 0)
                tf = cell.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = _ALIGN["left" if (i > 0 and j == 0) else align]
                _add_runs(p, val, self.pt(fs), col, bold)
        return tbl


class Ax:
    def __init__(self, s, left, right, bottom, top, xlim, ylim, xlog, ylog):
        self.s = s; self.L = left; self.R = right; self.B = bottom; self.T = top
        self.x0, self.x1 = xlim; self.y0, self.y1 = ylim
        self.xlog = xlog; self.ylog = ylog
        import math
        self._log = math.log10

    def _fx(self, dx):
        if self.xlog:
            t = (self._log(dx) - self._log(self.x0)) / (self._log(self.x1) - self._log(self.x0))
        else:
            t = (dx - self.x0) / (self.x1 - self.x0)
        return self.L + t * (self.R - self.L)

    def _fy(self, dy):
        if self.ylog:
            t = (self._log(dy) - self._log(self.y0)) / (self._log(self.y1) - self._log(self.y0))
        else:
            t = (dy - self.y0) / (self.y1 - self.y0)
        return self.B + t * (self.T - self.B)

    def xy(self, dx, dy):                      # -> mm top-origin
        return self.s.fig_xy(self._fx(dx), self._fy(dy))

    def wx(self, dw):                          # data x-width -> mm (linear axes)
        return abs(dw) / (self.x1 - self.x0) * (self.R - self.L) * self.s.FW

    def hy(self, dh):                          # data y-height -> mm (linear axes)
        return abs(dh) / (self.y1 - self.y0) * (self.T - self.B) * self.s.FH

    # ---- shapes ----------------------------------------------------------
    def _shape(self, kind, l, t, w, h):
        shp = self.s.slide.shapes.add_shape(kind, Mm(l), Mm(t), Mm(w), Mm(h))
        _no_shadow(shp)
        return shp

    def rect(self, x, y, w, h, fc, ec=INK, lw=0.5, rounded=False, round_adj=0.18,
             label=None, label_fs=7.0, label_color=INK, label_bold=False,
             label_italic=False, alpha=None):
        l, t = self.xy(x, y + h)
        shp = self._shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                          l, t, self.wx(w), self.hy(h))
        if fc is None:
            shp.fill.background()
        else:
            shp.fill.solid(); shp.fill.fore_color.rgb = C(fc)
        if ec is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = C(ec); shp.line.width = Pt(self.s.pt(lw))
        if rounded:
            try:
                shp.adjustments[0] = round_adj
            except Exception:
                pass
        if label is not None:
            tf = shp.text_frame
            _fill_tf(tf, label, self.s.pt(label_fs), label_color, label_bold,
                     label_italic, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, None)
        return shp

    def oval(self, cx, cy, r, fc, ec=None, lw=0.5, ry=None):
        ry = r if ry is None else ry
        l, t = self.xy(cx - r, cy + ry)
        shp = self._shape(MSO_SHAPE.OVAL, l, t, self.wx(2 * r), self.hy(2 * ry))
        if fc is None:
            shp.fill.background()
        else:
            shp.fill.solid(); shp.fill.fore_color.rgb = C(fc)
        if ec is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = C(ec); shp.line.width = Pt(self.s.pt(lw))
        return shp

    def line(self, x0, y0, x1, y1, color, lw=0.8, dashed=False, arrow=False):
        a = self.xy(x0, y0); b = self.xy(x1, y1)
        cn = self.s.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                               Mm(a[0]), Mm(a[1]), Mm(b[0]), Mm(b[1]))
        _no_shadow(cn)
        cn.line.color.rgb = C(color); cn.line.width = Pt(self.s.pt(lw))
        ln = cn.line._get_or_add_ln()
        if dashed:
            d = ln.makeelement(qn("a:prstDash"), {"val": "dash"}); ln.append(d)
        if arrow:
            e = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
            ln.append(e)
        return cn

    def arrow(self, x0, y0, x1, y1, color, lw=0.8, dashed=False):
        return self.line(x0, y0, x1, y1, color, lw, dashed, arrow=True)

    def freeform(self, pts, color, lw=1.2, closed=False, fill=None):
        """pts: list of (dx,dy) in data coords -> editable freeform line/shape."""
        mm = [self.xy(px, py) for px, py in pts]
        fb = self.s.slide.shapes.build_freeform(Mm(mm[0][0]), Mm(mm[0][1]), scale=1.0)
        fb.add_line_segments([(Mm(x), Mm(y)) for x, y in mm[1:]], close=closed)
        shp = fb.convert_to_shape()
        _no_shadow(shp)
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
        shp.line.color.rgb = C(color); shp.line.width = Pt(self.s.pt(lw))
        return shp

    def text(self, x, y, lines, fs, color, ha="center", va="center", bold=False,
             italic=False, w_local=None, h_local=None, line_spacing=None,
             wmm=None, hmm=None, rot=None):
        # Box is positioned in mm (anchored at the (x,y) data point per ha/va),
        # so this works on linear AND log axes. Use wmm/hmm for an explicit mm box
        # (required on log axes, where a data-unit width is not constant); w_local
        # gives a data-unit width on linear axes. rot rotates the box (e.g. 270 for
        # a vertical y-axis title).
        if wmm is not None:
            w_mm = wmm
        elif w_local is not None:
            w_mm = self.wx(w_local)
        else:
            w_mm = 30.0 if self.xlog else self.wx((self.x1 - self.x0) * 0.30)
        if hmm is not None:
            h_mm = hmm
        elif h_local is not None:
            h_mm = self.hy(h_local)
        else:
            h_mm = 6.0 if self.ylog else self.hy((self.y1 - self.y0) * 0.10)
        ax_mm, ay_mm = self.xy(x, y)            # anchor point, mm top-origin
        left = ax_mm - (w_mm / 2 if ha == "center" else w_mm if ha == "right" else 0)
        top = ay_mm - (h_mm / 2 if va == "center" else h_mm if va == "bottom" else 0)
        tb = self.s.slide.shapes.add_textbox(Mm(left), Mm(top), Mm(w_mm), Mm(h_mm))
        anchor = {"top": MSO_ANCHOR.TOP, "center": MSO_ANCHOR.MIDDLE,
                  "bottom": MSO_ANCHOR.BOTTOM}[va]
        _fill_tf(tb.text_frame, lines, self.s.pt(fs), color, bold, italic,
                 _ALIGN[ha], anchor, line_spacing)
        tb.text_frame.word_wrap = (w_local is not None or wmm is not None)
        if rot is not None:
            tb.rotation = rot
        return tb
